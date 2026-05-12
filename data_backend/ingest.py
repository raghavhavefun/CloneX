import hashlib
import json
import mimetypes
import uuid
import re
import os
import base64
import urllib.request
from datetime import datetime
from pathlib import Path

import chromadb
import numpy as np
import trafilatura
from chromadb.config import Settings as ChromaSettings
from pypdf import PdfReader
import pypdfium2 as pdfium
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from .settings import Settings
from .storage import MetadataStore


class IngestionService:
    EMBED_BATCH_SIZE = 64
    MAX_EXTRACT_TEXT_CHARS = int(os.getenv("MAX_EXTRACT_TEXT_CHARS", "400000"))
    MAX_EMBED_CHUNKS = int(os.getenv("MAX_EMBED_CHUNKS", "2500"))
    REJECTED_BINARY_EXTENSIONS = {
        ".exe", ".dll", ".so", ".dylib", ".bin", ".dat", ".iso",
        ".msi", ".apk", ".dmg", ".pkg", ".deb", ".rpm",
        ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz",
        ".mp3", ".wav", ".flac", ".mp4", ".mkv", ".avi", ".mov",
    }

    def __init__(self, settings: Settings, store: MetadataStore):
        self.settings = settings
        self.store = store
        self._embedder = None
        self._ocr = None
        self.chroma = chromadb.PersistentClient(
            path=str(settings.chroma_path),
            settings=ChromaSettings(anonymized_telemetry=False),
        )
        self.collection = self.chroma.get_or_create_collection(name="aria_memory")

    def _get_embedder(self):
        if self._embedder is None:
            from sentence_transformers import SentenceTransformer

            self._embedder = SentenceTransformer(self.settings.embed_model)
        return self._embedder

    def _get_ocr(self):
        if self._ocr is None:
            from paddleocr import PaddleOCR

            self._ocr = PaddleOCR(use_angle_cls=True, lang="en")
        return self._ocr

    def _run_ocr(self, image_input):
        """
        PaddleOCR API differs across versions.
        Try compatible call signatures and return best-effort result.
        """
        ocr = self._get_ocr()
        try:
            return ocr.ocr(image_input, cls=True)
        except TypeError:
            try:
                return ocr.ocr(image_input)
            except Exception:
                return []
        except Exception:
            return []

    def _chunk_text(self, text: str) -> list[str]:
        text = (text or "").strip()
        if not text:
            return []
        max_chars = int(self.settings.max_chunk_tokens * 4)
        overlap = int(self.settings.chunk_overlap * 4)
        chunks = []
        start = 0
        n = len(text)
        while start < n:
            end = min(n, start + max_chars)
            chunks.append(text[start:end])
            if end >= n:
                break
            start = max(0, end - overlap)
        return chunks

    def _safe_filename(self, name: str) -> str:
        base = (name or "file").strip().replace("\\", "_").replace("/", "_")
        base = re.sub(r"[^\w. ()-]+", "_", base)
        base = re.sub(r"\s+", " ", base).strip()
        return base[:180] if base else "file"

    def _owner_slug(self, owner_email: str | None) -> str:
        raw = (owner_email or "global").strip().lower()
        safe = re.sub(r"[^a-z0-9._-]+", "_", raw)
        return safe or "global"

    def _reject_binary_extension(self, original_name: str):
        suffix = Path(original_name or "").suffix.lower()
        if suffix in self.REJECTED_BINARY_EXTENSIONS:
            raise ValueError(f"Unsupported binary file type: {suffix}")

    def _extract_text_from_file(self, file_path: Path, mime: str) -> str:
        suffix = file_path.suffix.lower()
        if suffix in {".txt", ".md", ".json", ".py", ".csv", ".log"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".pdf":
            reader = PdfReader(str(file_path))
            text = "\n".join((p.extract_text() or "") for p in reader.pages)
            # OCR fallback for scanned/image PDFs where extract_text returns empty.
            if len((text or "").strip()) < 40:
                ocr_text = self._extract_text_from_pdf_via_ocr(file_path)
                if len((ocr_text or "").strip()) > len((text or "").strip()):
                    text = ocr_text
            return text
        if suffix == ".docx":
            doc = Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs)
        if suffix in {".pptx", ".ppt"}:
            prs = Presentation(str(file_path))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)  # type: ignore
            return "\n".join(lines)
        if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            wb = load_workbook(str(file_path), data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_txt = " | ".join("" if v is None else str(v) for v in row)
                    if row_txt.strip(" |"):
                        lines.append(row_txt)
            return "\n".join(lines)
        if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"}:
            result = self._run_ocr(str(file_path))
            lines = []
            if result:
                for block in result:
                    if not block:
                        continue
                    for item in block:
                        if item and len(item) > 1 and item[1]:
                            lines.append(str(item[1][0]))
            ocr_text = "\n".join(lines).strip()
            vlm_text = self._describe_image_with_vlm(file_path)
            if vlm_text and ocr_text:
                return f"[VLM_DENSE]\n{vlm_text}\n\n[OCR_TEXT]\n{ocr_text}"
            if vlm_text:
                return f"[VLM_DENSE]\n{vlm_text}"
            return ocr_text
        # Fallback raw text read
        return file_path.read_text(encoding="utf-8", errors="ignore")

    def _describe_image_with_vlm(self, file_path: Path) -> str:
        """
        Dense image understanding via local VLM (Ollama-compatible API).
        Non-fatal: returns empty string on any failure.
        """
        try:
            model = os.getenv("VLM_MODEL", "qwen3-vl:8b").strip() or "qwen3-vl:8b"
            base = os.getenv("OLLAMA_BASE", "http://127.0.0.1:11434").rstrip("/")
            url = f"{base}/api/chat"
            b64 = base64.b64encode(file_path.read_bytes()).decode("utf-8")
            payload = {
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": (
                            "Analyze this image densely and factually. Return rich details covering:\n"
                            "1) all visible objects/entities and attributes\n"
                            "2) scene layout and spatial relationships\n"
                            "3) any visible text/verbatim labels\n"
                            "4) actions/interactions/events\n"
                            "5) practical metadata for retrieval (keywords, likely domain/context)\n"
                            "Avoid speculation; mark uncertainty explicitly."
                        ),
                        "images": [b64],
                    }
                ],
                "stream": False,
            }
            body = json.dumps(payload).encode("utf-8")
            req = urllib.request.Request(
                url,
                data=body,
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=120) as resp:
                raw = resp.read().decode("utf-8", errors="ignore")
            data = json.loads(raw or "{}")
            text = ((data.get("message") or {}).get("content") or "").strip()
            return text[:20000]
        except Exception:
            return ""

    def _extract_text_from_pdf_via_ocr(self, file_path: Path, max_pages: int | None = None) -> str:
        lines = []
        try:
            pdf = pdfium.PdfDocument(str(file_path))
            page_count = len(pdf) if max_pages is None else min(len(pdf), max_pages)
            for page_idx in range(page_count):
                page = pdf[page_idx]
                bitmap = page.render(scale=2.0)  # type: ignore
                pil_img = bitmap.to_pil()
                result = self._run_ocr(np.array(pil_img))
                if not result:
                    continue
                for block in result:
                    if not block:
                        continue
                    for item in block:
                        if item and len(item) > 1 and item[1]:
                            txt = str(item[1][0]).strip()
                            if txt:
                                lines.append(txt)
        except Exception:
            return ""
        return "\n".join(lines)

    def _index_chunks(self, asset_id: str, chunks: list[str], metadata: dict):
        if not chunks:
            return
        if len(chunks) > int(self.MAX_EMBED_CHUNKS):
            chunks = chunks[: int(self.MAX_EMBED_CHUNKS)]
        embedder = self._get_embedder()
        batch_size = int(self.EMBED_BATCH_SIZE)
        for start in range(0, len(chunks), batch_size):
            end = min(start + batch_size, len(chunks))
            batch_chunks = chunks[start:end]
            batch_embeddings = embedder.encode(batch_chunks, normalize_embeddings=True).tolist()  # type: ignore
            ids = [f"{asset_id}:{idx}" for idx in range(start, end)]
            metadatas = []
            for idx in range(start, end):
                m = metadata.copy()
                m["chunk_index"] = idx
                m["asset_id"] = asset_id
                metadatas.append(m)
            self.collection.upsert(
                ids=ids,
                documents=batch_chunks,
                embeddings=batch_embeddings,
                metadatas=metadatas,
            )

    def ingest_file(self, file_bytes: bytes, original_name: str, source_type: str = "file", owner_email: str | None = None) -> dict:
        self._reject_binary_extension(original_name)
        asset_id = str(uuid.uuid4())
        created_at = datetime.utcnow().isoformat()
        hash_hex = hashlib.sha256(file_bytes).hexdigest()
        safe_original = self._safe_filename(original_name)
        safe_name = f"{asset_id}__{safe_original}"
        owner_slug = self._owner_slug(owner_email)
        raw_root = self.settings.data_vault_root / "raw" / owner_slug
        if (source_type or "").startswith("automation"):
            raw_root = raw_root / "automation"
        raw_root.mkdir(parents=True, exist_ok=True)
        raw_path = raw_root / safe_name
        raw_path.write_bytes(file_bytes)

        mime = mimetypes.guess_type(original_name)[0] or "application/octet-stream"
        asset = {
            "id": asset_id,
            "owner_email": (owner_email or "global").strip().lower(),
            "source_type": source_type,
            "name": original_name,
            "mime": mime,
            "path_or_url": str(raw_path),
            "size_bytes": len(file_bytes),
            "text_length": 0,
            "chunk_count": 0,
            "status": "ingested",
            "created_at": created_at,
        }

        text = self._extract_text_from_file(raw_path, mime)
        if len(text or "") > int(self.MAX_EXTRACT_TEXT_CHARS):
            text = (text or "")[: int(self.MAX_EXTRACT_TEXT_CHARS)]
        chunks = self._chunk_text(text)
        asset["text_length"] = len(text or "")
        asset["chunk_count"] = len(chunks)
        self.store.upsert_asset(asset)
        processed = {
            "asset_id": asset_id,
            "hash": hash_hex,
            "text_length": len(text or ""),
            "chunk_count": len(chunks),
            "text_preview": (text or "")[:2000],
            "created_at": created_at,
        }
        processed_root = self.settings.data_vault_root / "processed" / owner_slug
        processed_root.mkdir(parents=True, exist_ok=True)
        processed_path = processed_root / f"{asset_id}.json"
        processed_path.write_text(json.dumps(processed, indent=2), encoding="utf-8")
        extracted_text_path = processed_root / f"{asset_id}.txt"
        extracted_text_path.write_text(text or "", encoding="utf-8")

        self._index_chunks(
            asset_id,
            chunks,
            {
                "source_type": source_type,
                "mime": mime,
                "name": original_name,
                "owner_email": (owner_email or "global").strip().lower(),
            },
        )

        self.store.add_event(str(uuid.uuid4()), asset_id, "ingested", json.dumps({"chunk_count": len(chunks)}))
        return {**asset, "chunk_count": len(chunks)}

    def ingest_link(self, url: str, owner_email: str | None = None) -> dict:
        asset_id = str(uuid.uuid4())
        created_at = datetime.utcnow().isoformat()
        downloaded = trafilatura.fetch_url(url)
        text = trafilatura.extract(downloaded, include_links=True, include_formatting=False) if downloaded else ""
        if len(text or "") > int(self.MAX_EXTRACT_TEXT_CHARS):
            text = (text or "")[: int(self.MAX_EXTRACT_TEXT_CHARS)]

        asset = {
            "id": asset_id,
            "owner_email": (owner_email or "global").strip().lower(),
            "source_type": "link",
            "name": url,
            "mime": "text/url",
            "path_or_url": url,
            "size_bytes": len(text or ""),
            "text_length": len(text or ""),
            "chunk_count": 0,
            "status": "ingested",
            "created_at": created_at,
        }

        owner_slug = self._owner_slug(owner_email)
        processed_root = self.settings.data_vault_root / "processed" / owner_slug
        processed_root.mkdir(parents=True, exist_ok=True)
        processed_path = processed_root / f"{asset_id}.txt"
        processed_path.write_text(text or "", encoding="utf-8")

        chunks = self._chunk_text(text or "")
        asset["chunk_count"] = len(chunks)
        self.store.upsert_asset(asset)
        self._index_chunks(
            asset_id,
            chunks,
            {"source_type": "link", "mime": "text/url", "name": url, "owner_email": (owner_email or "global").strip().lower()},
        )

        self.store.add_event(str(uuid.uuid4()), asset_id, "ingested", json.dumps({"chunk_count": len(chunks)}))
        return {**asset, "chunk_count": len(chunks)}

    def ingest_text(self, text: str, title: str = "quick_note", owner_email: str | None = None) -> dict:
        asset_id = str(uuid.uuid4())
        created_at = datetime.utcnow().isoformat()
        cleaned = (text or "").strip()
        if not cleaned:
            cleaned = "[empty note]"
        if len(cleaned) > int(self.MAX_EXTRACT_TEXT_CHARS):
            cleaned = cleaned[: int(self.MAX_EXTRACT_TEXT_CHARS)]

        safe_title = self._safe_filename(title)
        safe_name = f"{asset_id}__{safe_title}.txt"
        owner_slug = self._owner_slug(owner_email)
        raw_root = self.settings.data_vault_root / "raw" / owner_slug
        raw_root.mkdir(parents=True, exist_ok=True)
        raw_path = raw_root / safe_name
        raw_path.write_text(cleaned, encoding="utf-8")

        asset = {
            "id": asset_id,
            "owner_email": (owner_email or "global").strip().lower(),
            "source_type": "text",
            "name": title,
            "mime": "text/plain",
            "path_or_url": str(raw_path),
            "size_bytes": len(cleaned.encode("utf-8")),
            "text_length": len(cleaned),
            "chunk_count": 0,
            "status": "ingested",
            "created_at": created_at,
        }

        processed_root = self.settings.data_vault_root / "processed" / owner_slug
        processed_root.mkdir(parents=True, exist_ok=True)
        processed_path = processed_root / f"{asset_id}.txt"
        processed_path.write_text(cleaned, encoding="utf-8")

        chunks = self._chunk_text(cleaned)
        asset["chunk_count"] = len(chunks)
        self.store.upsert_asset(asset)
        self._index_chunks(
            asset_id,
            chunks,
            {"source_type": "text", "mime": "text/plain", "name": title, "owner_email": (owner_email or "global").strip().lower()},
        )

        self.store.add_event(str(uuid.uuid4()), asset_id, "ingested", json.dumps({"chunk_count": len(chunks)}))
        return {**asset, "chunk_count": len(chunks)}

    def delete_asset(self, asset_id: str, owner_email: str | None = None) -> bool:
        asset = self.store.get_asset(asset_id, owner_email=owner_email)
        if not asset:
            return False

        ids = []
        for item in self.collection.get(where={"asset_id": asset_id}, include=[]).get("ids", []):
            if isinstance(item, list):
                ids.extend(item)
            else:
                ids.append(item)
        if ids:
            self.collection.delete(ids=ids)

        if asset["source_type"] == "file":
            try:
                p = Path(asset["path_or_url"])
                if p.exists():
                    p.unlink()
            except Exception:
                pass

        owner_slug = self._owner_slug(owner_email or asset.get("owner_email"))
        processed_path_json = self.settings.data_vault_root / "processed" / owner_slug / f"{asset_id}.json"
        processed_path_txt = self.settings.data_vault_root / "processed" / owner_slug / f"{asset_id}.txt"
        for p in [processed_path_json, processed_path_txt]:
            if p.exists():
                p.unlink()

        self.store.mark_deleted(asset_id)
        self.store.add_event(str(uuid.uuid4()), asset_id, "deleted", json.dumps({"ok": True}))
        return True

    def reprocess_asset(self, asset_id: str, owner_email: str | None = None) -> dict | None:
        asset = self.store.get_asset(asset_id, owner_email=owner_email)
        if not asset or asset.get("deleted_at"):
            return None

        source_type = asset.get("source_type", "")
        name = asset.get("name", "")
        self._reject_binary_extension(name)
        mime = asset.get("mime", "")
        path_or_url = asset.get("path_or_url", "")

        if source_type == "link":
            downloaded = trafilatura.fetch_url(path_or_url)
            text = trafilatura.extract(downloaded, include_links=True, include_formatting=False) if downloaded else ""
        else:
            p = Path(path_or_url)
            if not p.exists():
                return None
            text = self._extract_text_from_file(p, mime)
        if len(text or "") > int(self.MAX_EXTRACT_TEXT_CHARS):
            text = (text or "")[: int(self.MAX_EXTRACT_TEXT_CHARS)]

        chunks = self._chunk_text(text or "")
        self.collection.delete(where={"asset_id": asset_id})
        self._index_chunks(
            asset_id,
            chunks,
            {"source_type": source_type, "mime": mime, "name": name, "owner_email": (asset.get("owner_email") or owner_email or "global").strip().lower()},
        )

        asset["text_length"] = len(text or "")
        asset["chunk_count"] = len(chunks)
        asset["status"] = "ingested"
        self.store.upsert_asset(asset)

        processed_meta = {
            "asset_id": asset_id,
            "text_length": len(text or ""),
            "chunk_count": len(chunks),
            "text_preview": (text or "")[:2000],
            "reprocessed_at": datetime.utcnow().isoformat(),
        }
        owner_slug = self._owner_slug(owner_email or asset.get("owner_email"))
        processed_root = self.settings.data_vault_root / "processed" / owner_slug
        processed_root.mkdir(parents=True, exist_ok=True)
        processed_path = processed_root / f"{asset_id}.json"
        processed_path.write_text(json.dumps(processed_meta, indent=2), encoding="utf-8")
        processed_text_path = processed_root / f"{asset_id}.txt"
        processed_text_path.write_text(text or "", encoding="utf-8")

        self.store.add_event(str(uuid.uuid4()), asset_id, "reprocessed", json.dumps({"chunk_count": len(chunks)}))
        return {
            "id": asset_id,
            "name": name,
            "source_type": source_type,
            "mime": mime,
            "text_length": len(text or ""),
            "chunk_count": len(chunks),
            "status": "ingested",
        }
